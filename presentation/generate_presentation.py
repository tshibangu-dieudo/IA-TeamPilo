"""
Générateur de présentation PowerPoint pour TeamPilot AI
Conférence: AI FOR BUSINESS 2026
Université Lumière de Bujumbura (ULBU)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Palette de couleurs IBM
IBM_BLUE = RGBColor(15, 98, 254)  # Bleu IBM principal
IBM_DARK_BLUE = RGBColor(0, 67, 206)
IBM_LIGHT_BLUE = RGBColor(78, 168, 255)
IBM_GRAY = RGBColor(82, 95, 107)
IBM_LIGHT_GRAY = RGBColor(210, 213, 216)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

def create_presentation():
    """Crée la présentation complète"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Page de titre
    create_title_slide(prs)
    
    # Slide 2: Qui suis-je ?
    create_about_me_slide(prs)
    
    # Slide 3: Le problème
    create_problem_slide(prs)
    
    # Slide 4: Pourquoi l'IA ?
    create_why_ai_slide(prs)
    
    # Slide 5: Présentation TeamPilot AI
    create_teampilot_overview_slide(prs)
    
    # Slide 6: Fonctionnalités principales
    create_features_slide(prs)
    
    # Slide 7: Architecture IA
    create_ai_architecture_slide(prs)
    
    # Slide 8: Démonstration scénario
    create_demo_scenario_slide(prs)
    
    # Slide 9: Pourquoi IBM Granite ?
    create_why_granite_slide(prs)
    
    # Slide 10: Sécurité
    create_security_slide(prs)
    
    # Slide 11: Innovation
    create_innovation_slide(prs)
    
    # Slide 12: Impact
    create_impact_slide(prs)
    
    # Slide 13: Roadmap
    create_roadmap_slide(prs)
    
    # Slide 14: Démonstration Live
    create_demo_live_slide(prs)
    
    # Slide 15: Conclusion
    create_conclusion_slide(prs)
    
    # Sauvegarder
    output_path = "TeamPilot_AI_Conference_ULBU_2026.pptx"
    prs.save(output_path)
    print(f"✓ Présentation créée avec succès: {output_path}")
    return output_path

def add_title(slide, title_text):
    """Ajoute un titre stylisé à une slide"""
    # Ajouter un titre comme textbox en haut
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = IBM_BLUE
    p.alignment = PP_ALIGN.LEFT

def add_textbox(slide, left, top, width, height, text, font_size=18, 
                color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    """Ajoute une zone de texte stylisée"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), 
                                       Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return textbox

def add_bullet_points(slide, left, top, width, height, points, font_size=20):
    """Ajoute des points à puces"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), 
                                       Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = IBM_GRAY
        p.space_after = Pt(12)
    
    return textbox

def add_shape_box(slide, left, top, width, height, text, 
                  bg_color=IBM_BLUE, text_color=WHITE):
    """Ajoute une boîte colorée avec texte"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = bg_color
    
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    return shape

# ============================================================================
# SLIDE 1: PAGE DE TITRE
# ============================================================================
def create_title_slide(prs):
    """Slide 1: Page de titre élégante"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Fond bleu en haut
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(3)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = IBM_BLUE
    bg_shape.line.fill.background()
    
    # Titre principal
    add_textbox(slide, 1, 0.8, 8, 1, "TeamPilot AI", 
                font_size=60, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    # Sous-titre
    add_textbox(slide, 1, 1.6, 8, 0.6, 
                "Intelligence Artificielle pour la Gestion d'Équipe", 
                font_size=24, color=WHITE, align=PP_ALIGN.CENTER)
    
    # Slogan
    add_textbox(slide, 1, 3.5, 8, 0.8, 
                '"The AI proposes. The human decides."', 
                font_size=28, color=IBM_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    # Informations conférence
    add_textbox(slide, 1, 4.8, 8, 0.4, 
                "AI FOR BUSINESS 2026", 
                font_size=20, color=IBM_DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    add_textbox(slide, 1, 5.2, 8, 0.4, 
                "Université Lumière de Bujumbura (ULBU)", 
                font_size=18, color=IBM_GRAY, align=PP_ALIGN.CENTER)
    
    # Présentateur
    add_textbox(slide, 1, 6.2, 8, 0.3, 
                "Dieudo Tshibangu | Étudiant en Informatique de Gestion", 
                font_size=16, color=IBM_GRAY, align=PP_ALIGN.CENTER)
    
    # Technologies
    add_textbox(slide, 1, 6.7, 8, 0.3, 
                "IBM Granite • watsonx.ai • LangChain • Django • React", 
                font_size=14, color=IBM_LIGHT_GRAY, align=PP_ALIGN.CENTER)
    
    # Notes du présentateur
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """
[60 secondes]

Bonjour à toutes et à tous. Je suis Dieudo Tshibangu, étudiant en Informatique de Gestion à l'Université Lumière de Bujumbura.

Aujourd'hui, je suis honoré de vous présenter TeamPilot AI, une plateforme intelligente qui révolutionne la gestion d'équipe grâce à l'intelligence artificielle d'IBM Granite via watsonx.ai.

Notre slogan résume notre philosophie: "L'IA propose. L'humain décide."

TeamPilot AI n'est pas là pour remplacer les managers, mais pour les assister dans leurs décisions les plus critiques.

Ensemble, explorons comment l'IA peut transformer la gestion de projet en Afrique et dans le monde.
"""

# ============================================================================
# SLIDE 2: QUI SUIS-JE ?
# ============================================================================
def create_about_me_slide(prs):
    """Slide 2: Présentation personnelle"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_title(slide, "Qui suis-je ?")
    
    # Boîtes d'information
    add_shape_box(slide, 1, 2, 3.5, 1, "🎓 Étudiant\nInformatique de Gestion", 
                  IBM_BLUE, WHITE)
    add_shape_box(slide, 5.5, 2, 3.5, 1, "🇧🇮 Burundi\nULBU", 
                  IBM_DARK_BLUE, WHITE)
    
    # Points clés
    points = [
        "Passionné par l'IA et la gestion de projet",
        "Constat: Les équipes africaines manquent d'outils intelligents",
        "Mission: Démocratiser l'IA pour la gestion d'équipe",
        "Vision: Faire de TeamPilot AI la référence africaine"
    ]
    add_bullet_points(slide, 1.5, 3.5, 7, 3, points, font_size=22)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
[50 secondes]

Permettez-moi de me présenter rapidement. Je suis étudiant en Informatique de Gestion à l'ULBU, ici au Burundi.

Ma passion pour l'intelligence artificielle est née d'un constat simple: en Afrique, nos équipes utilisent des outils classiques qui ne répondent plus aux défis modernes.

J'ai développé TeamPilot AI avec une mission claire: démocratiser l'intelligence artificielle pour la gestion d'équipe, en la rendant accessible, abordable et efficace pour les PME africaines et les équipes universitaires.

Ma vision est de faire de TeamPilot AI la référence africaine en matière de gestion intelligente d'équipe.
"""

# ============================================================================
# SLIDE 3: LE PROBLÈME
# ============================================================================
def create_problem_slide(prs):
    """Slide 3: Le problème"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_title(slide, "Le Problème: Pourquoi les équipes échouent")
    
    # Statistiques choc
    add_textbox(slide, 1, 1.5, 8, 0.5, 
                "70% des projets échouent à cause de problèmes humains", 
                font_size=24, color=IBM_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    # Problèmes en colonnes
    problems = [
        ("🔴 Surcharge", "2.5"),
        ("⏰ Retards", "3.5"),
        ("🚫 Blocages", "4.5"),
        ("📊 Mauvaise visibilité", "5.5"),
        ("🤝 Communication", "2.5"),
        ("⚠️ Décisions tardives", "3.5")
    ]
    
    for i, (problem, left) in enumerate(problems):
        top = 2.8 if i < 3 else 4.3
        add_shape_box(slide, float(left), top, 1.8, 0.8, problem, 
                      IBM_LIGHT_GRAY, IBM_DARK_BLUE)
    
    # Conclusion
    add_textbox(slide, 1, 5.8, 8, 0.8, 
                "Les outils classiques (Excel, Trello) ne détectent pas ces problèmes", 
                font_size=20, color=IBM_GRAY, bold=True, align=PP_ALIGN.CENTER)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
[60 secondes]

Parlons du problème réel. Les statistiques sont alarmantes: 70% des projets échouent, non pas à cause de la technologie, mais à cause de problèmes humains.

La surcharge de travail: un développeur qui gère 150% de sa capacité finit par craquer.

Les retards en cascade: une tâche bloquée qui retarde tout le projet.

Le manque de visibilité: le manager découvre le problème trop tard.

La mauvaise communication: les équipes ne parlent pas le même langage.

Les décisions tardives: quand le manager réalise qu'il y a un problème, il est déjà trop tard.

Et le pire? Les outils classiques comme Excel, Trello ou Jira ne détectent pas ces problèmes. Ils se contentent d'afficher des tâches, sans analyser la situation.
"""

# ============================================================================
# SLIDE 4: POURQUOI L'IA ?
# ============================================================================
def create_why_ai_slide(prs):
    """Slide 4: Pourquoi l'IA ?"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_title(slide, "Pourquoi l'Intelligence Artificielle ?")
    
    # Avant/Après
    add_textbox(slide, 1, 2, 3.5, 0.6, "❌ AVANT (Outils classiques)", 
                font_size=22, color=IBM_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 5.5, 2, 3.5, 0.6, "✅ AVEC IA (TeamPilot)", 
                font_size=22, color=IBM_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    # Comparaisons
    comparisons = [
        ("Réactif", "Prédictif", 2.8),
        ("Manuel", "Automatique", 3.6),
        ("Décisions subjectives", "Data-driven", 4.4),
        ("Aucune prévision", "Anticipe les crises", 5.2)
    ]
    
    for before, after, top in comparisons:
        add_textbox(slide, 1, top, 3.5, 0.5, before, 
                    font_size=18, color=IBM_GRAY, align=PP_ALIGN.CENTER)
        add_textbox(slide, 5.5, top, 3.5, 0.5, after, 
                    font_size=18, color=IBM_DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
[50 secondes]

Pourquoi avons-nous besoin de l'intelligence artificielle ?

Les outils classiques sont réactifs. Vous découvrez le problème quand il est trop tard. L'IA est prédictive: elle anticipe les crises avant qu'elles n'arrivent.

Les outils classiques sont manuels. Vous devez tout calculer vous-même. L'IA automatise l'analyse: charges de travail, risques, dépendances.

Les outils classiques reposent sur des décisions subjectives. L'IA propose des décisions data-driven, basées sur des formules mathématiques et des règles métier validées.

Les outils classiques n'offrent aucune prévision. L'IA anticipe: "Dans 3 jours, ce projet sera en crise si vous ne redistribuez pas les tâches maintenant."

C'est ce changement de paradigme que TeamPilot AI apporte.
"""

# ============================================================================
# SLIDE 5: PRÉSENTATION TEAMPILOT AI
# ============================================================================
def create_teampilot_overview_slide(prs):
    """Slide 5: Vue d'ensemble TeamPilot AI"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_title(slide, "TeamPilot AI: Architecture Complète")
    
    # Schéma d'architecture en couches
    layers = [
        ("Frontend\nReact + Vite + Tailwind", 2, IBM_LIGHT_BLUE),
        ("Backend\nDjango REST Framework", 3.2, IBM_BLUE),
        ("Base de données\nPostgreSQL", 4.4, IBM_DARK_BLUE),
        ("Intelligence Artificielle\nIBM Granite via watsonx.ai", 5.6, IBM_BLUE)
    ]
    
    for layer, top, color in layers:
        add_shape_box(slide, 2, top, 6, 0.8, layer, color, WHITE)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
[50 secondes]

Voici l'architecture complète de TeamPilot AI.

En frontend, nous utilisons React avec Vite et Tailwind CSS pour une interface moderne et réactive.

Le backend repose sur Django REST Framework, un framework Python robuste et scalable, parfait pour les applications d'entreprise.

Les données sont stockées dans PostgreSQL, une base de données relationnelle performante et fiable.

Et au cœur du système: l'intelligence artificielle IBM Granite, accessible via watsonx.ai, intégrée avec LangChain.

Cette architecture garantit performance, sécurité et évolutivité.
"""

# ============================================================================
# SLIDE 6: FONCTIONNALITÉS PRINCIPALES
# ============================================================================
def create_features_slide(prs):
    """Slide 6: Fonctionnalités principales"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_title(slide, "Fonctionnalités Principales")
    
    # Grille de fonctionnalités 2x3
    features = [
        ("📁 Projets", "Gestion complète du cycle de vie", 1, 2.2),
        ("👥 Équipes", "Compétences et capacités", 4, 2.2),
        ("✅ Tâches", "Dépendances et statuts", 7, 2.2),
        ("📊 Workload", "Calcul en temps réel", 1, 4.2),
        ("⚠️ Risques", "Score composite IA", 4, 4.2),
        ("💬 Assistant", "Chat IBM Granite", 7, 4.2)
    ]
    
    for emoji_title, desc, left, top in features:
        add_shape_box(slide, left, top, 2.5, 1.5, 
                      f"{emoji_title}\n\n{desc}", IBM_LIGHT_BLUE, IBM_DARK_BLUE)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
[60 secondes]

TeamPilot AI offre six fonctionnalités principales intégrées.

Gestion de projets: cycle de vie complet, deadlines, statuts, ownership.

Gestion d'équipes: compétences, capacités hebdomadaires, disponibilités.

Gestion de tâches: avec dépendances, priorités, assignations intelligentes.

Calcul de workload en temps réel: formule mathématique précise qui détecte la surcharge avant qu'elle ne devienne critique.

Score de risque composite: quatre facteurs pondérés analysés par l'IA pour un score de 0 à 100%.

Assistant conversationnel IBM Granite: posez des questions en langage naturel, obtenez des réponses contextuelles basées sur vos données réelles.

Toutes ces fonctionnalités communiquent entre elles pour offrir une vue d'ensemble intelligente.
"""

# ============================================================================
# SLIDE 7: ARCHITECTURE IA
# ============================================================================
def create_ai_architecture_slide(prs):
    """Slide 7: Architecture IA détaillée"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_title(slide, "Architecture IA: Le Cerveau de TeamPilot")
    
    # Flux en cascade
    flow_steps = [
        ("1. Utilisateur", "Question / Demande", 1.8, IBM_LIGHT_BLUE),
        ("2. Intent Classifier", "Analyse l'intention", 2.6, IBM_BLUE),
        ("3. Router", "Dirige vers la bonne chaîne", 3.4, IBM_BLUE),
        ("4. Data Snapshot", "Collecte les données contextuelles", 4.2, IBM_DARK_BLUE),
        ("5. IBM Granite", "Génère la réponse intelligente", 5.0, IBM_BLUE),
        ("6. Validation", "Vérifie et retourne", 5.8, IBM_LIGHT_BLUE)
    ]
    
    for step, desc, top, color in flow_steps:
        add_shape_box(slide, 2, top, 6, 0.6, f"{step}: {desc}", color, WHITE)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
[70 secondes]

Voici le cerveau de TeamPilot AI: notre architecture d'intelligence artificielle.

Étape 1: L'utilisateur pose une question, par exemple "Qui est surchargé dans mon équipe ?"

Étape 2: L'Intent Classifier, notre premier modèle IA, analyse l'intention. Est-ce une question sur les personnes, les tâches, les risques ?

Étape 3: Le Router dirige la requête vers la bonne chaîne LangChain. Nous avons trois chaînes spécialisées: recommandations, explications de risque, et chat général.

Étape 4: Le Data Snapshot collecte toutes les données contextuelles pertinentes: membres de l'équipe, workloads, tâches, compétences, tout ce dont l'IA a besoin pour répondre intelligemment.

Étape 5: IBM Granite génère la réponse. Pas une réponse générique, mais une réponse basée sur VOS données, VOTRE contexte, VOTRE projet.

Étape 6: Validation et sécurité. Nous vérifions que la réponse respecte les permissions, ne contient pas d'hallucinations, et est cohérente.

Cette architecture garantit des réponses précises, sécurisées et contextuelles.
"""

# ============================================================================
# SLIDE 8: DÉMONSTRATION SCÉNARIO
# ============================================================================
def create_demo_scenario_slide(prs):
    """Slide 8: Scénario de démonstration"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_title(slide, "Scénario Réel: Le Manager face à la crise")
    
    # Scénario en étapes
    scenario = [
        ("🚨 Problème détecté", "David: 127.5% surchargé\nProjet CRM: 85.2% risque CRITIQUE", 2.2),
        ("🤖 IA suggère", "Réassigner 2 tâches de David vers Sarah\nSarah: seulement 32.5% de charge", 3.4),
        ("👨‍💼 Manager décide", "Alice (PM) valide la recommandation\nVoit la justification IA avant d'accepter", 4.6),
        ("✅ Action appliquée", "Tâches réassignées automatiquement\nNotifications envoyées", 5.8)
    ]
    
    for title, desc, top in scenario:
        add_shape_box(slide, 1.5, top, 7, 0.9, f"{title}\n{desc}", 
                      IBM_LIGHT_BLUE, IBM_DARK_BLUE)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
[60 secondes]

Voyons un scénario réel que nous allons démontrer.

Situation: David, développeur sur le projet CRM, est surchargé à 127.5%. Le projet CRM lui-même affiche un risque CRITIQUE de 85.2%. Toutes les tâches sont bloquées. C'est la crise.

L'IA détecte automatiquement ce problème. Elle analyse toute l'équipe et découvre que Sarah, sur le projet Mobile, n'est chargée qu'à 32.5%. Sarah a les compétences nécessaires.

L'IA suggère: "Réassignez 2 tâches de David vers Sarah. Cela réduira la charge de David et équilibrera l'équipe."

Le manager, Alice, voit cette recommandation. Elle lit la justification générée par IBM Granite. Elle comprend le raisonnement. Elle décide: "Oui, je valide."

L'action est appliquée instantanément. Les tâches sont réassignées. David et Sarah reçoivent des notifications. Le risque diminue.

L'IA propose. Le manager décide. C'est notre philosophie.
"""

# ============================================================================
# SLIDE 9: POURQUOI IBM GRANITE ?
# ============================================================================
def create_why_granite_slide(prs):
    """Slide 9: Pourquoi IBM Granite et watsonx"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_title(slide, "Pourquoi IBM Granite via watsonx.ai ?")
    
    # Raisons en boîtes
    reasons = [
        ("🔒 Sécurité\nEntreprise", 1, 2.2),
        ("🌍 Pas de\nVendor Lock", 3.2, 2.2),
        ("📊 Transparence\ndes données", 5.4, 2.2),
        ("🏢 Conçu pour\nle Business", 7.6, 2.2)
    ]
    
    for title, left, top in reasons:
        add_shape_box(slide, left, top, 2, 1.2, title, IBM_BLUE, WHITE)
    
    # Alternatives rejetées
    add_textbox(slide, 1, 4, 8, 0.5, 
                "❌ Pourquoi pas OpenAI/ChatGPT? Pas conçu pour l'entreprise, risques de confidentialité", 
                font_size=16, color=IBM_GRAY, align=PP_ALIGN.CENTER)
    
    add_textbox(slide, 1, 4.6, 8, 0.5, 
                "❌ Pourquoi pas un modèle local? Trop de ressources, maintenance complexe", 
                font_size=16, color=IBM_GRAY, align=PP_ALIGN.CENTER)
    
    add_textbox(slide, 1, 5.4, 8, 0.8, 
                "✅ IBM Granite: Le meilleur compromis performance/sécurité/coût", 
                font_size=20, color=IBM_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
[60 secondes]

Pourquoi avoir choisi IBM Granite via watsonx.ai plutôt qu'une autre solution ?

Quatre raisons principales.

Premièrement: la sécurité entreprise. IBM est un acteur de confiance depuis 100 ans. Les données de nos clients restent privées, conformes RGPD, jamais utilisées pour entraîner d'autres modèles.

Deuxièmement: pas de vendor lock-in. watsonx.ai supporte plusieurs modèles. Demain, si un meilleur modèle apparaît, nous pouvons switcher sans réécrire le code.

Troisièmement: la transparence. IBM Granite est entraîné sur des données d'entreprise ouvertes et documentées. Pas de boîte noire.

Quatrièmement: conçu pour le business. Granite comprend le langage métier, les KPIs, les processus d'entreprise.

Nous avons écarté OpenAI/ChatGPT: excellent pour le grand public, mais pas conçu pour les données sensibles d'entreprise.

Nous avons écarté les modèles locaux: trop lourds, trop coûteux à maintenir pour une startup.

IBM Granite est le meilleur compromis: performance, sécurité, coût raisonnable.
"""

# ============================================================================
# SLIDE 10: SÉCURITÉ
# ============================================================================
def create_security_slide(prs):
    """Slide 10: Sécurité et fiabilité"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_title(slide, "Sécurité: Zéro compromis")
    
    # Mesures de sécurité
    security_measures = [
        "🔐 Permissions granulaires par rôle (Admin, PM, Member, Executive)",
        "🛡️ Validation des réponses IA avant affichage",
        "🚫 Anti-hallucination: fallback templates si IA indisponible",
        "📝 Audit trail complet de toutes les actions",
        "🔒 Données chiffrées en transit et au repos",
        "✅ 222 tests automatisés pour garantir la cohérence"
    ]
    
    add_bullet_points(slide, 1.5, 2.5, 7, 4, security_measures, font_size=20)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
[50 secondes]

La sécurité est notre priorité absolue. Nous avons mis en place six niveaux de protection.

Premièrement: permissions granulaires. Chaque utilisateur ne voit que ce qu'il a le droit de voir. Un membre d'équipe ne peut pas accéder aux données d'une autre équipe.

Deuxièmement: validation des réponses IA. Toutes les sorties de Granite sont vérifiées avant d'être affichées à l'utilisateur.

Troisièmement: anti-hallucination. Si l'IA n'est pas disponible ou si elle produit une réponse douteuse, nous basculons automatiquement sur des templates déterministes.

Quatrièmement: audit trail complet. Chaque action est loggée: qui a fait quoi, quand, pourquoi.

Cinquièmement: chiffrement des données en transit et au repos. Vos données sont protégées à chaque étape.

Sixièmement: 222 tests automatisés qui s'exécutent à chaque modification du code, garantissant qu'aucune régression n'est introduite.

Avec TeamPilot AI, votre sécurité n'est pas négociable.
"""

# ============================================================================
# SLIDE 11: INNOVATION
# ============================================================================
def create_innovation_slide(prs):
    """Slide 11: Innovation et différenciation"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_title(slide, "Innovation: Ce qui nous rend uniques")
    
    # Tableau comparatif
    add_textbox(slide, 1, 2, 2, 0.5, "Fonctionnalité", 
                font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_shape_box(slide, 1, 2, 2, 0.5, "Fonctionnalité", IBM_BLUE, WHITE)
    
    competitors = ["Trello", "Asana", "Jira", "TeamPilot"]
    for i, comp in enumerate(competitors):
        color = IBM_BLUE if comp == "TeamPilot" else IBM_LIGHT_GRAY
        add_shape_box(slide, 3 + i*1.5, 2, 1.3, 0.5, comp, color, WHITE)
    
    features_comparison = [
        ("Gestion tâches", ["✅", "✅", "✅", "✅"], 2.7),
        ("Calcul workload", ["❌", "✅", "✅", "✅"], 3.4),
        ("Détection risques", ["❌", "❌", "⚠️", "✅"], 4.1),
        ("IA prédictive", ["❌", "❌", "❌", "✅"], 4.8),
        ("Assistant IA", ["❌", "❌", "❌", "✅"], 5.5)
    ]
    
    for feature, checks, top in features_comparison:
        add_textbox(slide, 1, top, 2, 0.5, feature, 
                    font_size=16, color=IBM_GRAY, align=PP_ALIGN.LEFT)
        for i, check in enumerate(checks):
            add_textbox(slide, 3 + i*1.5, top, 1.3, 0.5, check, 
                        font_size=18, align=PP_ALIGN.CENTER)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
[60 secondes]

Comparons TeamPilot AI avec les leaders du marché.

Gestion des tâches: tous les outils le font. C'est la base.

Calcul du workload: Trello ne le fait pas. Asana et Jira le font partiellement.

Détection des risques: ni Trello ni Asana. Jira a des dashboards, mais pas d'analyse prédictive.

IA prédictive: aucun d'entre eux. TeamPilot AI est le seul à anticiper les crises avant qu'elles n'arrivent.

Assistant IA conversationnel: TeamPilot AI uniquement. Vous pouvez littéralement demander "Qui est disponible pour m'aider sur ce projet ?" et obtenir une réponse contextuelle basée sur vos données réelles.

Notre innovation n'est pas d'ajouter de l'IA pour le buzz. Notre innovation est d'utiliser l'IA là où elle apporte une vraie valeur: prédire, anticiper, assister.
"""

# ============================================================================
# SLIDE 12: IMPACT
# ============================================================================
def create_impact_slide(prs):
    """Slide 12: Impact attendu"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_title(slide, "Impact: Qui peut bénéficier ?")
    
    # Cibles en colonnes
    targets = [
        ("🏢 Entreprises", "PME africaines\nÉquipes distribuées", 1.5, 2.5),
        ("🎓 Universités", "Projets étudiants\nRecherche collaborative", 4, 2.5),
        ("🚀 Startups", "Équipes agiles\nRessources limitées", 6.5, 2.5)
    ]
    
    for title, desc, left, top in targets:
        add_shape_box(slide, left, top, 2.2, 1.5, f"{title}\n\n{desc}", 
                      IBM_LIGHT_BLUE, IBM_DARK_BLUE)
    
    # Vision Afrique
    add_textbox(slide, 1, 4.5, 8, 1, 
                "Vision: Faire de l'Afrique un leader de l'IA appliquée à la gestion", 
                font_size=24, color=IBM_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    # Statistiques d'impact projetées
    add_textbox(slide, 1, 5.7, 8, 0.6, 
                "📊 Objectif: -30% de retards • +40% de productivité • -50% de burnout", 
                font_size=18, color=IBM_GRAY, align=PP_ALIGN.CENTER)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
[60 secondes]

Qui peut bénéficier de TeamPilot AI ?

Premièrement: les entreprises. Surtout les PME africaines qui n'ont pas les moyens d'embaucher des consultants en gestion de projet. TeamPilot AI leur donne accès à une intelligence artificielle de niveau entreprise à un coût abordable.

Deuxièmement: les universités. Les projets étudiants échouent souvent par manque de coordination. TeamPilot AI peut transformer la façon dont les étudiants collaborent sur des projets de recherche ou des travaux de groupe.

Troisièmement: les startups. Elles ont des ressources limitées, des équipes agiles, et des deadlines serrées. TeamPilot AI leur permet de rester organisés sans embaucher un Project Manager à temps plein.

Ma vision est plus large: faire de l'Afrique un leader de l'IA appliquée à la gestion. Nous avons les talents. Nous avons les problèmes à résoudre. Il nous manquait les outils. TeamPilot AI comble ce vide.

Nos objectifs d'impact: réduire les retards de 30%, augmenter la productivité de 40%, réduire le burnout de 50%.

C'est ambitieux. Mais c'est réalisable.
"""

# ============================================================================
# SLIDE 13: ROADMAP
# ============================================================================
def create_roadmap_slide(prs):
    """Slide 13: Roadmap future"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_title(slide, "Roadmap: Aujourd'hui et Demain")
    
    # Version actuelle
    add_textbox(slide, 1, 2, 4, 0.6, "✅ VERSION ACTUELLE (MVP)", 
                font_size=22, color=IBM_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    current_features = [
        "✓ Gestion complète projet/équipe/tâche",
        "✓ Calcul workload en temps réel",
        "✓ Score de risque composite",
        "✓ Recommandations IA (IBM Granite)",
        "✓ Assistant conversationnel",
        "✓ 222 tests automatisés"
    ]
    add_bullet_points(slide, 1.2, 2.8, 3.5, 2.5, current_features, font_size=16)
    
    # Version future
    add_textbox(slide, 5.5, 2, 4, 0.6, "🚀 VERSION FUTURE (Q3-Q4 2026)", 
                font_size=22, color=IBM_DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    future_features = [
        "📱 Application mobile (iOS/Android)",
        "📈 Analytics avancées avec ML",
        "🔮 Prédiction de vélocité par sprint",
        "🌍 Support multilingue complet",
        "🔗 Intégrations (Slack, Teams, GitHub)",
        "🤖 Auto-réassignation intelligente"
    ]
    add_bullet_points(slide, 5.7, 2.8, 3.5, 2.5, future_features, font_size=16)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
[60 secondes]

Parlons de notre roadmap.

Aujourd'hui, nous avons un MVP complet et fonctionnel. Gestion de projets, équipes, tâches. Calcul de workload en temps réel avec formules mathématiques précises. Score de risque composite pondéré. Recommandations générées par IBM Granite avec justifications. Assistant conversationnel qui comprend le contexte de vos projets. Et 222 tests automatisés qui garantissent la stabilité.

Ce MVP est prêt pour une utilisation en production.

Mais nous voyons plus loin. Pour Q3-Q4 2026, nous prévoyons:

Application mobile native pour iOS et Android, pour que les managers puissent gérer leurs équipes en déplacement.

Analytics avancées avec machine learning pour identifier des patterns dans vos projets.

Prédiction de vélocité: l'IA prédit combien de story points votre équipe pourra compléter au prochain sprint.

Support multilingue complet: français, anglais, swahili, kirundi, pour une adoption maximale en Afrique.

Intégrations avec Slack, Teams, GitHub, pour s'insérer dans vos workflows existants.

Et la fonctionnalité la plus ambitieuse: auto-réassignation intelligente. L'IA propose ET applique automatiquement les réassignations si vous l'autorisez.

C'est notre vision pour les 12 prochains mois.
"""

# ============================================================================
# SLIDE 14: DÉMONSTRATION LIVE
# ============================================================================
def create_demo_live_slide(prs):
    """Slide 14: Démonstration live"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    add_title(slide, "Démonstration Live")
    
    # Instructions pour le public
    add_textbox(slide, 1, 2.5, 8, 1, 
                "Voyons TeamPilot AI en action", 
                font_size=36, color=IBM_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    # Points à démontrer
    demo_points = [
        "1. Login en tant que Project Manager (demo_alice.pm)",
        "2. Dashboard: Voir le projet CRM en risque CRITIQUE (85.2%)",
        "3. Recommandations IA: Lire la justification IBM Granite",
        "4. Accepter une recommandation: Réassignation automatique",
        "5. Chat Assistant: 'Qui est surchargé dans mon équipe ?'",
        "6. Vérification: Voir les workloads mis à jour en temps réel"
    ]
    add_bullet_points(slide, 1.5, 3.8, 7, 3, demo_points, font_size=20)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
[180 secondes - Démonstration pratique]

Maintenant, le moment que vous attendiez: la démonstration live.

[Ouvrir le navigateur]

Je vais me connecter en tant que Project Manager Alice. Voici le dashboard.

[Montrer le dashboard]

Vous voyez ici le projet CRM Platform Rewrite. Le score de risque est affiché en rouge: 85.2%, niveau CRITIQUE. C'est l'alerte maximale.

[Cliquer sur le projet]

Entrons dans les détails. Vous voyez David Chen, notre développeur principal. Sa barre de workload est rouge: 127.5%, critically overloaded. Il a deux tâches critiques bloquées depuis plusieurs jours.

[Cliquer sur Recommandations]

L'IA a détecté ce problème. Regardez: elle propose de réassigner une tâche de David vers Sarah Kim, qui n'est chargée qu'à 32.5%.

[Lire la justification]

La justification générée par IBM Granite explique pourquoi: Sarah a les compétences React nécessaires, elle est sous-utilisée, et cette réassignation réduirait le risque du projet de 15 points.

[Accepter la recommandation]

Je clique sur "Accepter". En une seconde, la tâche est réassignée. David et Sarah reçoivent des notifications.

[Montrer le chat]

Maintenant, utilisons l'assistant IA. Je tape: "Qui est surchargé dans mon équipe ?"

[Attendre la réponse]

Granite analyse toutes mes données et répond: "David Chen est critically overloaded à 127.5%. Alice est à 110%, Marie à 110%, Lucas à 100%."

C'est une réponse basée sur MES données, en temps réel.

[Retour au dashboard]

Et vous voyez, le workload de David a déjà baissé. Le système a recalculé automatiquement.

Voilà TeamPilot AI en action. Des questions sur ce que vous venez de voir ?
"""

# ============================================================================
# SLIDE 15: CONCLUSION
# ============================================================================
def create_conclusion_slide(prs):
    """Slide 15: Conclusion et remerciements"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Fond dégradé bleu (simulé avec une grande boîte)
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = IBM_LIGHT_BLUE
    bg_shape.line.fill.background()
    
    add_title(slide, "TeamPilot AI: L'avenir de la gestion d'équipe")
    
    # Messages clés
    add_textbox(slide, 1, 2.5, 8, 0.6, 
                "L'IA propose. L'humain décide.", 
                font_size=32, color=IBM_DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    # Résumé en 3 points
    summary = [
        "✓ Détection prédictive des crises projet",
        "✓ IBM Granite pour des décisions data-driven",
        "✓ Conçu pour l'Afrique, utilisable partout"
    ]
    add_bullet_points(slide, 2, 3.5, 6, 1.5, summary, font_size=22)
    
    # Remerciements
    add_textbox(slide, 1, 5.2, 8, 0.5, 
                "Merci à l'ULBU et à tous les participants", 
                font_size=24, color=IBM_DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
    
    # Contact
    add_textbox(slide, 1, 6, 8, 0.4, 
                "📧 Contact: dieudo.tshibangu@teampilot.ai", 
                font_size=18, color=IBM_GRAY, align=PP_ALIGN.CENTER)
    
    # Questions
    add_textbox(slide, 1, 6.6, 8, 0.5, 
                "❓ VOS QUESTIONS", 
                font_size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = """
[50 secondes]

Pour conclure.

TeamPilot AI, c'est trois choses essentielles:

Premièrement: la détection prédictive. Nous n'attendons pas que la crise arrive. Nous l'anticipons et vous donnons les outils pour l'éviter.

Deuxièmement: IBM Granite pour des décisions data-driven. Pas de feeling, pas d'intuition. Des données, des formules, des recommandations justifiées.

Troisièmement: conçu pour l'Afrique, utilisable partout. Nous comprenons les défis spécifiques des PME africaines, mais notre solution est universelle.

Notre slogan résume tout: "L'IA propose. L'humain décide."

TeamPilot AI n'est pas là pour remplacer le manager. Il est là pour donner au manager les super-pouvoirs dont il a besoin pour réussir.

Je remercie l'Université Lumière de Bujumbura pour cette opportunité, et je remercie chacun d'entre vous pour votre attention.

Je suis maintenant prêt à répondre à toutes vos questions.

Merci.
"""

# ============================================================================
# MAIN EXECUTION
# ============================================================================
if __name__ == "__main__":
    print("🚀 Génération de la présentation TeamPilot AI...")
    print("=" * 60)
    output_file = create_presentation()
    print("=" * 60)
    print(f"✅ Présentation générée: {output_file}")
    print("\n📋 Contenu:")
    print("   • 15 slides professionnelles")
    print("   • Notes du présentateur pour chaque slide")
    print("   • Design IBM (bleu/blanc)")
    print("   • Prête pour AI FOR BUSINESS 2026")
    print("\n🎯 Prochaines étapes:")
    print("   1. Ouvrir le fichier .pptx")
    print("   2. Réviser les notes du présentateur")
    print("   3. Ajouter des images si nécessaire")
    print("   4. Pratiquer la présentation (15-18 min)")
